from __future__ import annotations

from datetime import datetime
from pathlib import Path

from app.utils.file_utils import ensure_dir, write_utf8

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - handled at runtime
    Presentation = None


PPT_EXTENSIONS: set[str] = {".pptx", ".ppt"}


PPT_GUIDANCE = ".ppt 파일은 직접 처리하기 어렵습니다. PowerPoint에서 .pptx로 저장한 뒤 다시 시도하세요."


def _list_ppt_files(input_dir: Path) -> list[Path]:
    if not input_dir.exists():
        return []
    return sorted(
        [path for path in input_dir.iterdir() if path.is_file() and path.suffix.lower() in PPT_EXTENSIONS],
        key=lambda path: path.name,
    )


def _safe_text(text: str) -> str:
    return "\n".join(line.rstrip() for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")).strip()


def _paragraph_lines(text_frame: object, indent: str = "") -> list[str]:
    lines: list[str] = []
    paragraphs = getattr(text_frame, "paragraphs", [])
    for paragraph in paragraphs:
        text = _safe_text("".join(run.text for run in getattr(paragraph, "runs", [])) or getattr(paragraph, "text", ""))
        if not text:
            continue
        level = int(getattr(paragraph, "level", 0) or 0)
        prefix = "  " * max(level, 0)
        lines.append(f"{indent}{prefix}- {text}" if text else "")
    return [line for line in lines if line.strip()]


def _extract_shape_lines(shape: object) -> list[str]:
    lines: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text_frame = shape.text_frame
        lines.extend(_paragraph_lines(text_frame))
    elif getattr(shape, "has_table", False):
        table = shape.table
        for row in table.rows:
            row_text = [
                _safe_text(getattr(cell, "text", "")) or ""
                for cell in row.cells
            ]
            row_text = [cell for cell in row_text if cell]
            if row_text:
                lines.append("- " + " | ".join(row_text))
    return lines


def _extract_slide_title(slide: object, body_lines: list[str]) -> tuple[str, list[str]]:
    title = ""
    remaining: list[str] = []

    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = _safe_text(shape.text_frame.text)
            if not text:
                continue
            if not title and getattr(shape, "is_placeholder", False):
                placeholder = getattr(shape, "placeholder_format", None)
                placeholder_type = str(getattr(placeholder, "type", "")) if placeholder else ""
                if "TITLE" in placeholder_type.upper() or "CENTER_TITLE" in placeholder_type.upper():
                    title = text
                    continue
            remaining.extend(_paragraph_lines(shape.text_frame))
        elif getattr(shape, "has_table", False):
            remaining.extend(_extract_shape_lines(shape))

    if not title:
        for idx, line in enumerate(remaining):
            if line.startswith("- "):
                candidate = line[2:].strip()
                if candidate:
                    title = candidate
                    remaining = remaining[:idx] + remaining[idx + 1 :]
                    break

    if not title:
        title = "텍스트 없음"

    body_lines.extend(remaining)
    return title, body_lines


def _extract_notes_text(slide: object) -> str:
    try:
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        text = _safe_text(notes_text_frame.text)
        return text
    except Exception:
        return ""


def _append_error_log(error_log_path: Path, messages: list[str]) -> Path:
    if not messages:
        return error_log_path
    ensure_dir(error_log_path.parent)
    timestamp = datetime.now().isoformat(timespec="seconds")
    with error_log_path.open("a", encoding="utf-8") as f:
        for message in messages:
            f.write(f"[{timestamp}] {message}\n")
    return error_log_path


def _render_ppt_markdown(
    original_name: str,
    slide_count: int,
    slides: list[dict[str, object]],
) -> str:
    lines: list[str] = [
        "# PPT 교육자료 추출 결과",
        "",
        f"- 원본 파일: {original_name}",
        f"- 슬라이드 수: {slide_count}",
        "",
    ]

    for item in slides:
        slide_number = int(item.get("slide_number", 0) or 0)
        title = str(item.get("title", "텍스트 없음") or "텍스트 없음")
        body_lines = [str(line) for line in item.get("body_lines", []) if str(line).strip()]
        notes_text = str(item.get("notes_text", "") or "").strip()

        lines.append(f"## Slide {slide_number}")
        lines.append("")
        lines.append(title)
        lines.append("")

        if body_lines:
            lines.append("본문")
            lines.extend(body_lines)
            lines.append("")
        else:
            lines.append("텍스트 없음")
            lines.append("")

        if notes_text:
            lines.append("발표자 노트")
            lines.append("")
            lines.append(notes_text)
            lines.append("")

    return "\n".join(lines).strip() + "\n"


def extract_ppt_files(input_dir: Path, output_dir: Path, force: bool = False) -> dict[str, list[Path]]:
    ppt_files = _list_ppt_files(input_dir)
    extracts_dir = ensure_dir(output_dir / "ppt_extracts")

    if not ppt_files:
        print("PPT 대상 파일이 없습니다.")
        return {"ppt_extracts": []}

    if Presentation is None:
        raise RuntimeError("python-pptx가 설치되어 있지 않습니다. pip install python-pptx 를 실행하세요.")

    extract_paths: list[Path] = []
    error_messages: list[str] = []

    for ppt_path in ppt_files:
        output_md_path = extracts_dir / f"{ppt_path.stem}.ppt.md"

        if not force and output_md_path.exists():
            extract_paths.append(output_md_path)
            continue

        if ppt_path.suffix.lower() == ".ppt":
            error_messages.append(f"{ppt_path.name}: {PPT_GUIDANCE}")
            print(PPT_GUIDANCE)
            continue

        try:
            prs = Presentation(str(ppt_path))
            slides_output: list[dict[str, object]] = []
            for slide_index, slide in enumerate(prs.slides, start=1):
                body_lines: list[str] = []
                title, body_lines = _extract_slide_title(slide, body_lines)
                notes_text = _extract_notes_text(slide)
                slides_output.append(
                    {
                        "slide_number": slide_index,
                        "title": title,
                        "body_lines": body_lines,
                        "notes_text": notes_text,
                    }
                )

            markdown_text = _render_ppt_markdown(
                original_name=ppt_path.name,
                slide_count=len(prs.slides),
                slides=slides_output,
            )
            write_utf8(output_md_path, markdown_text)
            extract_paths.append(output_md_path)
        except Exception as exc:  # pragma: no cover - runtime extraction errors
            error_message = f"{ppt_path.name}: {exc}"
            error_messages.append(error_message)
            print(error_message)

    if error_messages:
        error_log_path = _append_error_log(extracts_dir / "ppt_errors.log", error_messages)
        return {"ppt_extracts": extract_paths, "errors": [error_log_path]}

    return {"ppt_extracts": extract_paths}
